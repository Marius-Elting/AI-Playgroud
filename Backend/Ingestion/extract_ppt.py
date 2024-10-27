from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from dotenv import load_dotenv
from openai import OpenAI


load_dotenv()

OPENAI_API_KEY= os.getenv("OPENAI_API_KEY")

def extract_text_and_images(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    text_content = []
    image_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = s.image
                        image_bytes = image.blob
                        image_filename = f'image_{image_count}.jpg'
                        image_path = os.path.join(output_dir, image_filename)
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        image_count += 1
                    elif s.has_text_frame:
                        text_content.append(s.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_filename = f'image_{image_count}.jpg'
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                image_count += 1
            elif shape.has_text_frame:
                text_content.append(shape.text)
    
    text_filename = os.path.join(output_dir, 'extracted_text.txt')
    with open(text_filename, 'w') as f:
        f.write('\n'.join(text_content))



def translate(text):
    client = OpenAI()
    prompt = f"Translate the following text to english, please only return the text 1:1 translated to english: '{text}'"

    completion = client.chat.completions.create(
            messages= [
                {"role": "system", "content":"You are a helpful assistant that can help people with anything"},
                {"role": "user", "content": prompt}
            ],
            model="gpt-4o-mini",
            stream=False
        )
    if completion.choices[0].message.content == None:
        return ""
    return completion.choices[0].message.content

    result = openai_service.call_openai_without_streaming(prompt)

    # Placeholder for translation logic
    return result # Example: reversing the text as a dummy translation

def translate_pptx_text(pptx_path, output_pptx_path):
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if s.has_text_frame:
                        original_text = s.text
                        translated_text = translate(original_text)
                        s.text = translated_text
            elif shape.has_text_frame:
                original_text = shape.text
                translated_text = translate(original_text)
                shape.text = translated_text
    
    prs.save(output_pptx_path)

if __name__ == "__main__":
    file_path = './test.pptx'
    current_dir = os.path.dirname(os.path.abspath(__file__))
    absolute_path = os.path.join(current_dir, file_path)
    output_dir = 'test_output'
    extract_text_and_images(absolute_path, output_dir)
    
    output_pptx_path = os.path.join(current_dir, 'translated_test.pptx')
    translate_pptx_text(absolute_path, output_pptx_path)